from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from scripts.ingestion.models import AbstractIngestor, UnsupportedFileError
import logging
from scripts.utils.image_utils import get_project_image_dir, infer_project_root, record_image_metadata

logger = logging.getLogger("pptx_ingestor")


# def _infer_project_root(file_path: Path) -> Path:
#     """
#     Extracts the base project directory from a known data/projects/{project_name}/... path.
#     Returns: Path to data/projects/{project_name}
#     """
#     parts = file_path.resolve().parts
#     for i in range(len(parts) - 2):
#         if parts[i] == "data" and parts[i + 1] == "projects":
#             return Path(*parts[:i + 3])
#     return file_path.parent


class PptxIngestor(AbstractIngestor):
    """
    Ingestor for PPTX files.
    """

    def ingest(self, filepath: str) -> list[tuple[str, dict]]:
        if not filepath.endswith(".pptx"):
            raise UnsupportedFileError("File is not a .pptx file.")

        extracted_data = []
        try:

            prs = Presentation(filepath)
            file_path = Path(filepath)

            # 🔍 Infer project root and name
            project_root = infer_project_root(file_path)
            project_name = project_root.name              # ✅ CORRECT
            image_dir = get_project_image_dir(project_name)

            logger.info(f"[PPTX] Ingesting file: {file_path}")
            logger.info(f"[PPTX] Image output folder: {image_dir}")

            file_stem = file_path.stem

            for i, slide in enumerate(prs.slides):
                slide_number = i + 1
                text_on_slide = []
                image_counter = 0

                for shape in slide.shapes:
                    logger.debug(f"[PPTX] Slide {slide_number}, Shape type: {shape.shape_type}")
                    if hasattr(shape, "name"):
                        logger.debug(f"[PPTX] Shape name: {shape.name}")

                    if hasattr(shape, "text") and shape.text:
                        text_on_slide.append(shape.text.strip())

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_counter += 1
                        img_name = f"{file_stem}_slide{slide_number}_img{image_counter}.png"
                        rel_dir = Path("input") / "cache" / "images"
                        out_path = project_root / rel_dir / img_name
                        out_path.parent.mkdir(parents=True, exist_ok=True)

                        try:
                            with open(out_path, "wb") as f:
                                f.write(shape.image.blob)
                            logger.info(f"[PPTX] Saved image to: {out_path}")
                        except Exception as img_err:
                            logger.warning(f"[PPTX] Failed to save image on slide {slide_number}: {img_err}")
                            continue

                        try:
                            out_path = out_path.resolve()
                            rel_to_input = out_path.relative_to(project_root / "input")
                            img_rel = str(rel_to_input)
                        except Exception as e:
                            logger.warning(f"[PPTX] Failed to compute relative image path: {e}")
                            # img_rel = str(rel_dir / img_name)  # fallback

                        # Track images per slide
                        slide_meta = {
                            "slide_number": slide_number,
                            "type": "slide_content",
                            "doc_type": "pptx",
                        }
                        text_content = "\n".join(text_on_slide).strip()
                        has_images = False

                        for shape in slide.shapes:
                            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                image_counter += 1
                                img_name = f"{file_stem}_slide{slide_number}_img{image_counter}.png"
                                out_path = image_dir / img_name

                                try:
                                    with open(out_path, "wb") as f:
                                        f.write(shape.image.blob)

                                    logger.info(f"[PPTX] Saved image to: {out_path}")
                                    record_image_metadata(slide_meta, out_path, project_root)
                                    has_images = True

                                except Exception as img_err:
                                    logger.warning(f"[PPTX] Failed to save image on slide {slide_number}: {img_err}")
                                    continue

                        # Emit slide chunk only once
                        if text_content or has_images:
                            extracted_data.append((text_content or "[Image-only slide]", slide_meta))
                # Presenter notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        notes_meta = {
                            "slide_number": slide_number,
                            "type": "presenter_notes",
                            "doc_type": "pptx"
                        }
                        formatted_notes = (
                            f"Presenter Notes (Slide {slide_number}):\n"
                            f"{notes_text}"
                        )
                        extracted_data.append((formatted_notes, notes_meta))

                # Fallback: no image, but text exists
                if text_on_slide and image_counter == 0:
                    slide_content = "\n".join(text_on_slide).strip()
                    if slide_content:
                        slide_meta = {
                            "slide_number": slide_number,
                            "type": "slide_content",
                            "doc_type": "pptx"
                        }
                        extracted_data.append((slide_content, slide_meta))

        except Exception as e:
            logger.error(f"[PPTX] Fatal error processing file {filepath}: {e}", exc_info=True)
            raise UnsupportedFileError(f"Error processing PPTX file {filepath}: {e}")

        return extracted_data
